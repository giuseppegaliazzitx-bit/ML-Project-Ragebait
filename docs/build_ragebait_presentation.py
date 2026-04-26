from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "ragebait_detection_presentation.pptx"
OUTLINE = ROOT / "docs" / "ragebait_detection_presentation_outline.md"


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(28, 42, 66)
BLUE = RGBColor(47, 103, 246)
RED = RGBColor(214, 68, 68)
GREEN = RGBColor(55, 145, 100)
GOLD = RGBColor(180, 130, 35)
GRAY = RGBColor(95, 104, 120)
LIGHT_GRAY = RGBColor(242, 244, 248)
MID_GRAY = RGBColor(220, 224, 230)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(20, 24, 32)


def set_font(run, size=20, bold=False, color=BLACK, font="Aptos"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, text, x, y, w, h, size=20, bold=False, color=BLACK,
                align=PP_ALIGN.LEFT, font="Aptos", margin=0.04, auto_size=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    if auto_size:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color, font=font)
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, title, 0.72, 0.36, 8.9, 0.48, size=24, bold=True, color=NAVY)
    if subtitle:
        add_textbox(slide, subtitle, 0.75, 0.88, 7.8, 0.28, size=9.5, color=GRAY)
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(1.18), Inches(1.06), Inches(0.035)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()


def add_footer(slide, idx):
    add_textbox(slide, "ML Ragebait Detection Project", 0.72, 7.13, 3.4, 0.18, size=7.5, color=GRAY)
    add_textbox(slide, str(idx), 12.25, 7.13, 0.36, 0.18, size=7.5, color=GRAY, align=PP_ALIGN.RIGHT)


def add_bullets(slide, bullets, x, y, w, h, size=15.5, color=BLACK, gap=0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(gap)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def add_metric_card(slide, label, value, note, x, y, w=2.7, h=1.18, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = MID_GRAY
    shape.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_textbox(slide, label, x + 0.18, y + 0.14, w - 0.32, 0.22, size=8.8, bold=True, color=GRAY)
    add_textbox(slide, value, x + 0.18, y + 0.42, w - 0.32, 0.34, size=22, bold=True, color=NAVY)
    add_textbox(slide, note, x + 0.18, y + 0.82, w - 0.32, 0.22, size=8.5, color=GRAY)


def add_table(slide, headers, rows, x, y, w, h, font_size=9.5, header_fill=NAVY,
              first_col_bold=False, col_widths=None):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if col_widths:
        for i, frac in enumerate(col_widths):
            table.columns[i].width = int(Inches(w) * frac)
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            for run in p.runs:
                set_font(run, size=font_size, bold=True, color=WHITE)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.RIGHT if c > 0 and str(val).replace(".", "", 1).replace("%", "").isdigit() else PP_ALIGN.LEFT
                for run in p.runs:
                    set_font(run, size=font_size, bold=(first_col_bold and c == 0), color=BLACK)
    return table_shape


def add_bar_chart(slide, labels, series, x, y, w, h, title=None, ymax=1.0):
    if title:
        add_textbox(slide, title, x, y - 0.28, w, 0.22, size=9.5, bold=True, color=GRAY)
    chart_left = Inches(x + 0.45)
    chart_top = Inches(y)
    chart_w = Inches(w - 0.62)
    chart_h = Inches(h - 0.55)
    axis = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, chart_left, chart_top, chart_w, chart_h)
    axis.fill.background()
    axis.line.color.rgb = MID_GRAY
    n = len(labels)
    group_w = chart_w / n
    colors = [GOLD, BLUE, GREEN, RED]
    for s_idx, (name, vals) in enumerate(series):
        bar_w = group_w * 0.26
        offset = (s_idx - (len(series) - 1) / 2) * bar_w * 1.15
        for i, val in enumerate(vals):
            bh = int(chart_h * (val / ymax))
            bx = int(chart_left + group_w * i + group_w * 0.5 + offset - bar_w / 2)
            by = int(chart_top + chart_h - bh)
            bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, bx, by, int(bar_w), bh)
            bar.fill.solid()
            bar.fill.fore_color.rgb = colors[s_idx % len(colors)]
            bar.line.fill.background()
            add_textbox(slide, f"{val:.2f}", bx / 914400, by / 914400 - 0.20, bar_w / 914400 + 0.18, 0.16,
                        size=6.5, color=GRAY, align=PP_ALIGN.CENTER, margin=0, auto_size=False)
    for i, label in enumerate(labels):
        lx = (chart_left + group_w * i + group_w * 0.04) / 914400
        add_textbox(slide, label, lx, y + h - 0.42, group_w / 914400 * 0.92, 0.34, size=7, color=GRAY,
                    align=PP_ALIGN.CENTER, margin=0)
    legend_x = x + 0.46
    for s_idx, (name, _) in enumerate(series):
        swatch = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(legend_x), Inches(y + h - 0.06), Inches(0.12), Inches(0.12))
        swatch.fill.solid()
        swatch.fill.fore_color.rgb = colors[s_idx % len(colors)]
        swatch.line.fill.background()
        add_textbox(slide, name, legend_x + 0.17, y + h - 0.11, 1.5, 0.22, size=7.5, color=GRAY, margin=0)
        legend_x += 1.7


def add_image(slide, path, x, y, w=None, h=None):
    kwargs = {}
    if w:
        kwargs["width"] = Inches(w)
    if h:
        kwargs["height"] = Inches(h)
    return slide.shapes.add_picture(str(ROOT / path), Inches(x), Inches(y), **kwargs)


def blank_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    add_title(slide, title, subtitle)
    add_footer(slide, len(prs.slides))
    return slide


def make_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = "Machine Learning Ragebait Detection Project"
    prs.core_properties.subject = "Final project presentation"
    prs.core_properties.author = "Group 8"

    # 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_textbox(slide, "Machine Learning Ragebait Detection", 0.78, 1.05, 8.8, 0.74, size=31, bold=True, color=NAVY)
    add_textbox(slide, "A supervised NLP system for detecting abusive and rage-inducing social-media language", 0.82, 1.88, 8.5, 0.34, size=14, color=GRAY)
    add_textbox(slide, "Group 8 - Edison Cheah, Giuseppe Galiazzi, Lawson Herger, Andrew Lin", 0.82, 2.36, 8.8, 0.28, size=10.5, color=GRAY)
    add_metric_card(slide, "FINAL BINARY TEST", "0.9229", "BERT F1", 0.82, 3.35, accent=BLUE)
    add_metric_card(slide, "FINAL MULTICLASS TEST", "0.6405", "BERT macro F1", 3.75, 3.35, accent=RED)
    add_metric_card(slide, "GOLD-LABEL DATASET", "12,490", "human-labeled posts", 6.68, 3.35, accent=GREEN)
    add_bullets(slide, [
        "Final benchmark uses human labels, frozen splits, and a three-tier model ladder.",
        "The central finding is methodological: weak-label agreement was not enough; gold-label evaluation changed the scientific story."
    ], 0.82, 5.08, 9.2, 0.8, size=13.5, color=BLACK)
    add_footer(slide, 1)

    # 2
    slide = blank_slide(prs, "Research Problem", "Why this task is more than ordinary toxicity detection")
    add_bullets(slide, [
        "Short social-media posts compress sarcasm, abuse, provocation, threats, and ordinary profanity into very little context.",
        "A model must separate surface offensiveness from intent: trolling, derogation, hate speech, and normal speech can share vocabulary.",
        "The practical objective is reliable classification of rage-inducing or abusive language without treating every controversial post as ragebait."
    ], 0.9, 1.55, 5.65, 2.65, size=16)
    add_table(slide, ["Boundary", "Why it is difficult"], [
        ["Profanity vs. Hate Speech", "Explicit words can mask stronger harm or threat intent."],
        ["Trolling vs. Derogatory", "Taunting style overlaps with direct personal insult."],
        ["Normal vs. Trolling", "Controversy and sarcasm are not automatically abuse."]
    ], 6.85, 1.55, 5.45, 2.1, font_size=9.7, col_widths=[0.36, 0.64])
    add_metric_card(slide, "HARDER TARGET", "5 classes", "Normal, Profanity, Trolling, Derogatory, Hate Speech", 6.85, 4.35, 5.45, 1.12, accent=RED)

    # 3
    slide = blank_slide(prs, "Project Arc", "The project intentionally preserves both the first attempt and the final system")
    add_table(slide, ["Phase", "Supervision", "Key output", "Main lesson"], [
        ["Iteration 1", "LLM weak labels", "507,682 labeled posts; balanced 32k training file", "Scores measured agreement with Qwen, not human judgment."],
        ["Iteration 2", "Human labels", "12,490-row gold-label benchmark", "Final claims are grounded in held-out human annotation."]
    ], 0.85, 1.55, 11.7, 1.45, font_size=9.4, first_col_bold=True, col_widths=[0.15, 0.16, 0.36, 0.33])
    add_bullets(slide, [
        "The weak-label pipeline was valuable engineering practice: import, label, filter, balance, train, and evaluate.",
        "The scientific pivot was necessary because high performance against model-generated labels does not prove real validity.",
        "The final deck therefore emphasizes Iteration 2 while using Iteration 1 as motivation."
    ], 1.05, 3.55, 10.75, 1.45, size=15)
    add_metric_card(slide, "ARCHIVED WEAK-LABEL BERT", "0.8735", "macro F1 against LLM labels", 1.05, 5.35, 3.4, 1.12, accent=GOLD)
    add_metric_card(slide, "FINAL GOLD-LABEL BERT", "0.9229", "binary F1 against human labels", 4.72, 5.35, 3.4, 1.12, accent=BLUE)
    add_metric_card(slide, "FINAL MULTICLASS BERT", "0.6405", "macro F1 on 5 classes", 8.39, 5.35, 3.4, 1.12, accent=RED)

    # 4
    slide = blank_slide(prs, "Final Dataset", "Human-labeled abusive-language corpus used for all final experiments")
    add_table(slide, ["Label", "Rows", "Share"], [
        ["Normal", "5,053", "40.46%"],
        ["Profanity", "1,582", "12.67%"],
        ["Trolling", "4,537", "36.32%"],
        ["Derogatory", "862", "6.90%"],
        ["Hate Speech", "456", "3.65%"]
    ], 0.92, 1.45, 4.55, 2.45, font_size=10.6, first_col_bold=True, col_widths=[0.48, 0.25, 0.27])
    add_bullets(slide, [
        "Binary task maps Normal to 0 and every abusive category to 1.",
        "Multiclass task keeps all five labels: Normal, Profanity, Trolling, Derogatory, Hate Speech.",
        "Class imbalance is substantial, especially for Hate Speech and Derogatory."
    ], 6.05, 1.55, 5.75, 1.6, size=14.5)
    add_metric_card(slide, "BINARY BALANCE", "59.54%", "ragebait / abusive examples", 6.05, 3.7, 2.8, 1.1, accent=RED)
    add_metric_card(slide, "MINORITY CLASS", "3.65%", "Hate Speech examples", 9.05, 3.7, 2.8, 1.1, accent=GOLD)
    add_textbox(slide, "Dataset source: trolldata.csv, used as the final gold-label benchmark.", 0.94, 4.65, 5.6, 0.3, size=9.5, color=GRAY)

    # 5
    slide = blank_slide(prs, "Frozen Splits", "Same train, validation, and test rows across all final models")
    add_table(slide, ["Split", "Rows", "Normal", "Ragebait / abusive"], [
        ["Train", "9,992", "4,042", "5,950"],
        ["Validation", "1,249", "506", "743"],
        ["Test", "1,249", "505", "744"]
    ], 0.85, 1.55, 5.8, 1.65, font_size=11, first_col_bold=True, col_widths=[0.28, 0.22, 0.25, 0.25])
    add_table(slide, ["Split", "Normal", "Profanity", "Trolling", "Derogatory", "Hate Speech"], [
        ["Train", "4,042", "1,265", "3,630", "690", "365"],
        ["Validation", "506", "158", "454", "86", "45"],
        ["Test", "505", "159", "453", "86", "46"]
    ], 0.85, 4.05, 11.6, 1.45, font_size=9.8, first_col_bold=True, col_widths=[0.18, 0.16, 0.16, 0.16, 0.17, 0.17])
    add_bullets(slide, [
        "Split policy: stratified 80/10/10 with seed 42.",
        "Model selection used validation performance; final claims use held-out test results.",
        "Freezing the split prevents accidental improvements from split drift."
    ], 7.0, 1.55, 5.1, 1.8, size=13.7)

    # 6
    slide = blank_slide(prs, "Modeling Ladder", "Three tiers test whether richer representations improve performance")
    add_table(slide, ["Tier", "Model family", "Representation", "Role"], [
        ["1", "Logistic Regression, Linear SVC", "TF-IDF unigrams and bigrams", "Strong lexical baselines"],
        ["2", "PyTorch FFNN", "Train-only embeddings with mean pooling", "Compact neural baseline"],
        ["3", "BERT", "bert-base-uncased contextual encoder", "Fine-grained intent modeling"]
    ], 0.85, 1.55, 11.6, 1.8, font_size=10, first_col_bold=True, col_widths=[0.08, 0.28, 0.32, 0.32])
    add_bullets(slide, [
        "TF-IDF establishes whether labels are mostly lexical.",
        "FFNN tests whether learned dense interactions help beyond linear baselines.",
        "BERT tests whether context helps separate overlapping abuse categories."
    ], 1.05, 3.95, 5.2, 1.35, size=14.5)
    add_table(slide, ["BERT setting", "Binary", "Multiclass"], [
        ["Max length", "80", "128"],
        ["Epochs", "2", "4"],
        ["Selection metric", "Validation F1", "Validation macro F1"]
    ], 7.0, 3.85, 4.8, 1.45, font_size=10.2, col_widths=[0.42, 0.29, 0.29])

    # 7
    slide = blank_slide(prs, "Implementation", "Final workspace is reproducible and artifact-driven")
    add_bullets(slide, [
        "iteration2/src/data creates label maps and canonical splits.",
        "iteration2/src/training runs TF-IDF baselines, FFNN, and BERT.",
        "iteration2/src/evaluation saves metrics, confusion matrices, and hard-error CSV files.",
        "Outputs include summaries, tokenizer snapshots, checkpoints, training histories, and split manifests."
    ], 0.95, 1.55, 5.65, 2.55, size=14.2)
    code = (
        "weight_i = N / (K * count_i)\n"
        "\n"
        "embeddings = self.embedding(input_ids)\n"
        "mask = (input_ids != self.pad_index).unsqueeze(-1)\n"
        "pooled = (embeddings * mask).sum(dim=1) / mask.sum(dim=1).clamp(min=1)\n"
        "logits = self.classifier(pooled)"
    )
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.55), Inches(5.25), Inches(2.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.color.rgb = MID_GRAY
    add_textbox(slide, code, 7.2, 1.78, 4.85, 1.55, size=9.2, color=NAVY, font="Courier New")
    add_textbox(slide, "Two implementation details: exact class weights for imbalance and a simple mean-pooled FFNN baseline.", 7.0, 4.05, 5.2, 0.55, size=11.5, color=GRAY)

    # 8
    slide = blank_slide(prs, "Iteration 1 Results", "A useful pipeline, but an insufficient endpoint")
    add_table(slide, ["Model / run", "Accuracy", "Positive F1", "Macro F1", "Notes"], [
        ["Initial BERT", "0.4055", "0.5770", "0.2885", "Collapsed toward positive predictions"],
        ["Raw-text Linear SVC", "0.8744", "0.8484", "0.8706", "Strongest exact-split classical baseline"],
        ["Tuned BERT", "0.8767", "0.8533", "0.8735", "Best archived weak-label result"]
    ], 0.85, 1.55, 11.65, 1.62, font_size=9.5, first_col_bold=True, col_widths=[0.26, 0.14, 0.14, 0.14, 0.32])
    add_bullets(slide, [
        "The tuned BERT barely exceeded the strongest classical baseline.",
        "More importantly, the evaluation target was still model-generated labels.",
        "This became the reason to pivot to a human-labeled benchmark."
    ], 1.0, 3.75, 6.2, 1.35, size=15)
    add_metric_card(slide, "WEAK-LABELED CORPUS", "507,682", "LLM-labeled examples", 7.8, 3.75, 3.9, 1.1, accent=GOLD)
    add_metric_card(slide, "HIGH-CONFIDENCE POOL", "102,104", "confidence >= 0.95", 7.8, 5.05, 3.9, 1.1, accent=BLUE)

    # 9
    slide = blank_slide(prs, "Experiment 1: Binary Classification", "Held-out test results for Normal vs. ragebait / abusive")
    add_table(slide, ["Model", "Tier", "Accuracy", "Precision", "Recall", "F1"], [
        ["Logistic Regression", "1", "0.8551", "0.8371", "0.9395", "0.8854"],
        ["Linear SVC", "1", "0.8655", "0.8780", "0.8992", "0.8884"],
        ["FFNN", "2", "0.8671", "0.8734", "0.9086", "0.8906"],
        ["BERT", "3", "0.9079", "0.9210", "0.9247", "0.9229"]
    ], 0.85, 1.45, 7.0, 2.2, font_size=10.3, first_col_bold=True, col_widths=[0.34, 0.11, 0.14, 0.14, 0.13, 0.14])
    add_metric_card(slide, "BEST MODEL", "BERT", "Tier 3 contextual encoder", 8.3, 1.55, 3.35, 1.1, accent=BLUE)
    add_metric_card(slide, "F1 GAIN", "+0.0322", "over FFNN baseline", 8.3, 2.9, 3.35, 1.1, accent=GREEN)
    add_bullets(slide, [
        "The binary task is already strong for lexical models, so BERT's gain is meaningful rather than automatic.",
        "The FFNN only marginally improves over Linear SVC, suggesting the coarse boundary is mostly lexical.",
        "BERT still gives the best balance of precision and recall on held-out data."
    ], 0.98, 4.35, 10.8, 1.35, size=14)

    # 10
    slide = blank_slide(prs, "Binary Confusion Matrix", "BERT reduces both false positives and false negatives on the final benchmark")
    add_image(slide, "iteration2/outputs/exp1_binary_bert/bert_confusion_matrix.png", 0.95, 1.4, w=5.45)
    add_bullets(slide, [
        "BERT held-out test: 0.9079 accuracy and 0.9229 F1.",
        "False negatives remain important because abusive posts missed by the classifier carry moderation risk.",
        "False positives matter because controversial but non-abusive text should not be over-flagged."
    ], 7.0, 1.75, 4.8, 2.0, size=14.2)
    add_metric_card(slide, "PRECISION", "0.9210", "positive predictions usually correct", 7.0, 4.45, 2.35, 1.05, accent=BLUE)
    add_metric_card(slide, "RECALL", "0.9247", "most abusive posts found", 9.65, 4.45, 2.35, 1.05, accent=RED)

    # 11
    slide = blank_slide(prs, "Experiment 2: Five-Class Classification", "Held-out test results for the full label set")
    add_table(slide, ["Model", "Tier", "Accuracy", "Micro F1", "Macro F1"], [
        ["Logistic Regression", "1", "0.6645", "0.6645", "0.5524"],
        ["Linear SVC", "1", "0.6797", "0.6797", "0.5403"],
        ["FFNN", "2", "0.6373", "0.6373", "0.5101"],
        ["BERT", "3", "0.7390", "0.7390", "0.6405"]
    ], 0.85, 1.45, 6.55, 2.1, font_size=10.6, first_col_bold=True, col_widths=[0.4, 0.12, 0.16, 0.16, 0.16])
    add_image(slide, "iteration2/outputs/exp2_multiclass_bert/bert_test_confusion_matrix.png", 8.05, 1.35, w=3.95)
    add_bullets(slide, [
        "The five-class task exposes category-boundary ambiguity that binary detection hides.",
        "BERT improves test macro F1 by 0.0881 over the best held-out baseline.",
        "The remaining errors concentrate around Trolling, Derogatory, and Hate Speech."
    ], 0.98, 4.25, 6.5, 1.35, size=14)

    # 12
    slide = blank_slide(prs, "Class-wise Performance", "BERT improves every class against the best non-transformer baseline")
    labels = ["Normal", "Profanity", "Trolling", "Derog.", "Hate"]
    add_bar_chart(slide, labels, [
        ("Logistic Regression", [0.8215, 0.6108, 0.6202, 0.3398, 0.3696]),
        ("BERT", [0.8772, 0.7349, 0.6853, 0.4384, 0.4667]),
    ], 0.8, 1.55, 7.1, 4.0, ymax=0.95)
    add_table(slide, ["Class", "LogReg F1", "BERT F1", "Gain"], [
        ["Normal", "0.8215", "0.8772", "+0.0557"],
        ["Profanity", "0.6108", "0.7349", "+0.1241"],
        ["Trolling", "0.6202", "0.6853", "+0.0651"],
        ["Derogatory", "0.3398", "0.4384", "+0.0986"],
        ["Hate Speech", "0.3696", "0.4667", "+0.0971"]
    ], 8.15, 1.55, 4.35, 2.35, font_size=9.8, first_col_bold=True, col_widths=[0.36, 0.21, 0.21, 0.22])
    add_bullets(slide, [
        "Largest gains appear where surface vocabulary alone is least reliable: Profanity, Derogatory, and Hate Speech."
    ], 8.22, 4.45, 3.95, 0.65, size=12.2, color=GRAY)

    # 13
    slide = blank_slide(prs, "Error Analysis", "The hardest failures are semantic and boundary-driven")
    add_table(slide, ["Confusion pair", "Symmetric errors"], [
        ["Derogatory <-> Trolling", "93"],
        ["Normal <-> Trolling", "85"],
        ["Profanity <-> Trolling", "66"],
        ["Normal <-> Profanity", "33"],
        ["Derogatory <-> Hate Speech", "29"]
    ], 0.85, 1.45, 5.3, 2.3, font_size=10.2, first_col_bold=True, col_widths=[0.72, 0.28])
    add_bullets(slide, [
        "Profanity can hide more severe hostile intent.",
        "Trolling and derogatory abuse often share direct insult vocabulary.",
        "Long rants mix sarcasm, complaint, threat, and provocation in one post.",
        "OOV rates do not explain the hard BERT errors: hard-error OOV is 7.61%, close to the full test OOV rate of 7.83%."
    ], 6.75, 1.45, 5.4, 2.15, size=13.4)
    add_metric_card(slide, "TOP DIRECTIONAL ERROR", "68", "Trolling predicted as Derogatory", 0.85, 4.45, 3.25, 1.1, accent=RED)
    add_metric_card(slide, "HARD-ERROR OOV", "7.61%", "not a vocabulary-hole story", 4.35, 4.45, 3.25, 1.1, accent=GREEN)
    add_metric_card(slide, "MAIN TAKEAWAY", "semantic", "not just lexical", 7.85, 4.45, 3.25, 1.1, accent=BLUE)

    # 14
    slide = blank_slide(prs, "Compute Tradeoffs", "Performance gains must be weighed against training and inference cost")
    add_table(slide, ["Model", "Binary train", "Binary predict", "Multiclass train", "Multiclass predict"], [
        ["Logistic Regression", "0.047 s", "0.00039 s", "16.314 s", "0.00097 s"],
        ["Linear SVC", "0.032 s", "0.00045 s", "0.471 s", "0.00146 s"],
        ["FFNN", "25.388 s", "0.01973 s", "87.860 s", "0.1501 s"],
        ["BERT", "4893.595 s", "64.869 s", "376.786 s", "9.860 s"]
    ], 0.85, 1.45, 11.65, 2.05, font_size=9.8, first_col_bold=True, col_widths=[0.28, 0.18, 0.18, 0.18, 0.18])
    add_bullets(slide, [
        "TF-IDF models are effectively free and remain valuable baselines.",
        "BERT is the only model with a serious compute cost, but GPU training makes the multiclass run practical.",
        "Operationally, BERT is most justified when category distinction matters, not merely when coarse binary screening is enough."
    ], 1.05, 4.1, 10.9, 1.35, size=14)
    add_metric_card(slide, "BINARY BERT", "81.6 min", "CPU training run", 1.05, 5.8, 3.1, 0.85, accent=GOLD)
    add_metric_card(slide, "MULTICLASS BERT", "6.3 min", "CUDA training run", 4.55, 5.8, 3.1, 0.85, accent=GREEN)
    add_metric_card(slide, "PRACTICAL USE", "batch", "offline moderation scoring", 8.05, 5.8, 3.1, 0.85, accent=BLUE)

    # 15
    slide = blank_slide(prs, "Conclusion", "Final results and next steps")
    add_metric_card(slide, "FINAL BINARY BERT", "0.9229", "held-out test F1", 0.95, 1.55, 3.35, 1.15, accent=BLUE)
    add_metric_card(slide, "FINAL MULTICLASS BERT", "0.6405", "held-out test macro F1", 4.65, 1.55, 3.35, 1.15, accent=RED)
    add_metric_card(slide, "BEST BASELINE GAP", "+0.0881", "multiclass macro F1", 8.35, 1.55, 3.35, 1.15, accent=GREEN)
    add_bullets(slide, [
        "Classical lexical baselines are strong and should always be reported.",
        "The FFNN adds little on binary detection and underperforms the best linear model on multiclass macro F1.",
        "BERT is the only model that consistently improves aggregate scores and class-wise robustness.",
        "The remaining research problem is better handling of overlapping abusive categories, not simply more parameters."
    ], 1.05, 3.35, 10.7, 1.95, size=15)
    add_table(slide, ["Promising next step", "Reason"], [
        ["Grouped splits by author/source", "Reduce leakage and test real generalization."],
        ["More context per post", "Resolve sarcasm, threat, and reply-chain ambiguity."],
        ["Sharper annotation guidance", "Clarify Trolling vs. Derogatory vs. Hate Speech."]
    ], 1.05, 5.65, 10.9, 1.0, font_size=9.4, col_widths=[0.36, 0.64])

    # 16
    slide = blank_slide(prs, "References", "Core sources used by the project")
    add_bullets(slide, [
        "Trawling for Trolling: A Dataset. arXiv:2008.00525.",
        "Devlin, Chang, Lee, and Toutanova. BERT: Pre-training of Deep Bidirectional Transformers for Language Understanding. NAACL-HLT, 2019.",
        "Pedregosa et al. Scikit-learn: Machine Learning in Python. Journal of Machine Learning Research, 2011.",
        "Project repository: github.com/giuseppegaliazzitx-bit/ML-Project-Ragebait"
    ], 1.0, 1.65, 10.9, 2.3, size=14)
    add_textbox(slide, "End", 1.0, 5.2, 2.0, 0.44, size=24, bold=True, color=NAVY)

    prs.save(OUT)


def write_outline():
    outline = """# Ragebait Detection Presentation

Generated deck: `docs/ragebait_detection_presentation.pptx`

## Slide List

1. Machine Learning Ragebait Detection
2. Research Problem
3. Project Arc
4. Final Dataset
5. Frozen Splits
6. Modeling Ladder
7. Implementation
8. Iteration 1 Results
9. Experiment 1: Binary Classification
10. Binary Confusion Matrix
11. Experiment 2: Five-Class Classification
12. Class-wise Performance
13. Error Analysis
14. Compute Tradeoffs
15. Conclusion
16. References

## Speaker Emphasis

- The strongest story is the pivot from LLM weak-label agreement to human-labeled evaluation.
- Binary classification confirms BERT is best, but also shows lexical baselines are strong.
- Multiclass classification is the most scientifically interesting task because it exposes semantic boundary ambiguity.
- The remaining challenge is not OOV vocabulary; it is category overlap among trolling, derogation, profanity, and hate speech.
"""
    OUTLINE.write_text(outline, encoding="utf-8")


if __name__ == "__main__":
    make_deck()
    write_outline()
    print(f"Wrote {OUT}")
    print(f"Wrote {OUTLINE}")
